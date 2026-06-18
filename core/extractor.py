from pypdf import PdfReader
from docx import Document
from pptx import Presentation



def extract_pdf(uploaded_file):
   try:
     text = ""
     reader = PdfReader(uploaded_file)
     for page in reader.pages:
         page_text = page.extract_text()
         if page_text:
             text += page_text + "\n"
     return text
   except Exception as e:
        raise ValueError(f"PDF Extraction Failed: {e}")
        

#Function to Exctract text from Docx files
def extract_docx(uploaded_file):
   try:
    text = ""
    doc = Document(uploaded_file)
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text
   except Exception as e:
        raise ValueError(f"DOCX Extraction Failed: {e}")
        

#Function to extract text from PPTX files
def extract_pptx(uploaded_file):
    try:
        text = ""
        prs = Presentation(uploaded_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        raise ValueError(f"PPTX Extraction Failed: {e}")
        



def extract_txt(uploaded_file):
    return uploaded_file.read().decode("utf-8")

def extract_text(uploaded_file, filename):
    
    filename = filename.lower()

    if filename.endswith(".pdf"):
        return extract_pdf(uploaded_file)

    elif filename.endswith(".docx"):
        return extract_docx(uploaded_file)

    elif filename.endswith(".pptx"):
        return extract_pptx(uploaded_file)

    elif filename.endswith(".txt"):
        return extract_txt(uploaded_file)

    raise ValueError("Unsupported file type")