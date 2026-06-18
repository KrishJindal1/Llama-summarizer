import streamlit as st
import ollama
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
# Added dictionary to store the temperature and instructions for each tone option
tone_config = {
    "Professional": {
        
        "instruction": "Use polished, professional, business-oriented language."
    },
    "Formal": {
        
        "instruction": "Use formal vocabulary, complete sentences, and a serious tone."
    },
    "Casual": {
        
        "instruction": "Use conversational, friendly, and approachable language."
    },
    "Executive": {
        
        "instruction": "Focus on key insights, decisions, outcomes, and business impact."
    },
    "Technical": {
        
        "instruction": "Use precise technical terminology and maintain technical accuracy."
    },
    "Academic": {
        
        "instruction": "Use objective, analytical, and scholarly language."
    },
    "Marketing": {
       
        "instruction": "Make the content engaging, persuasive, and audience-focused."
    },
    "Simple English": {
        
        "instruction": "Use simple words, short sentences, and explain concepts clearly."
    }
}
   


# Function to analyze the document, generate summary and rewrite based on the selected tone and rewrite target
def analyze_document(document_text,selected_model,temperature,tone,rewrite_target,target_length):
    with st.status("📄 Processing Document...", expanded=True) as status:

        st.write("Generating summary...")
        
        # Call the Ollama API to get the summary or rewrite based on the user's choice
        try:
            summary_response = ollama.chat(
                model=selected_model,
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert document analyst."
                    },
                    {
                        "role": "user",
                        #prompt for better summarized content.
                        "content": f"""
                            Generate a summary.

                                Target Lengths:
                                    {target_length} 

                                Requirements:
                                - Capture all key points.
                                - Preserve important facts.
                                - Do not omit critical information.
                                - Use clear headings where appropriate.
                                - Do not leave the summary unfinished.
                                
                                Document:
                                {document_text}
                                """
                    }
                ],
                options={
                    "temperature": 1,
                    "num_predict": 2000 # set a high token limit to avoid the summary being cut off due to token limits.
                }
            )
        except Exception as e:
            st.error(f"Summary Generation Failed: {e}")
            return
        summary = summary_response["message"]["content"]
        st.success("Summary generated")

        st.write("Generating rewritten content...")
        config = tone_config[tone]
        instruction = config["instruction"]

        
        try:#Api call to rewrite the content based on the selected tone and rewrite target 
            rewrite_response = ollama.chat(
                model=selected_model,
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert content rewriter."
                    },
                    {
                        "role": "user",
                        "content": f"""
                            You are an expert editor.

                            Task:
                                Rewrite the content using a {tone} tone.

                            Tone Guidelines:
                                {instruction}

                                Requirements:
                                - Preserve all important information.
                                - Do not change factual content.
                                - Improve readability and organization.
                                - Keep the rewritten version approximately the same length.
                                - Use headings or bullet points when appropriate.
                                - Do not leave the rewritten version unfinished.
                                

                            Content:
                            {summary if rewrite_target == "Summary" else document_text}
                            """
                    }
                ],
                options={
                    "temperature": temperature,
                }
            )
        except Exception as e:
            st.error(f"Rewriting Failed: {e}")
            return
        rewritten_text = rewrite_response["message"]["content"]
        st.success("Rewritten text generated")

        status.update(
            label="✅ Analysis Complete",
            state="complete",
        )

    #Display the summary and rewritten text in separate tabs
    sum_tab1, rewrite_tab2 = st.tabs(["📄 Summary", "✍️ Rewritten"])
    with sum_tab1:
            st.subheader("📄 Summary")
            st.write(summary)
    with rewrite_tab2:
            st.subheader("✍️ Rewritten Text")
            st.write(rewritten_text)

# Function to extract text from PDF files
def extract_pdf(uploaded_file):
   try:
     text = ""
     reader = PdfReader(uploaded_file)
     for page in reader.pages:
         page_text = page.extract_text()
         if page_text:
             text += page_text + "\n"
     return text
   except Exception as e:
        st.error(f"PDF Extraction Failed: {e}")
        return ""

#Function to Exctract text from Docx files
def extract_docx(uploaded_file):
   try:
    text = ""
    doc = Document(uploaded_file)
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text
   except Exception as e:
        st.error(f"DOCX Extraction Failed: {e}")
        return ""

#Function to extract text from PPTX files
def extract_pptx(uploaded_file):
    try:
        text = ""
        prs = Presentation(uploaded_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(
            f"PPTX Extraction Failed: {e}")
        return ""

# Function to show the preview of the extracted text from the uploaded document
def show_preview(document_text):
   try:
    st.text_area(
        "Preview",
        document_text[:2000],
        height=200
    )
   except Exception as e:
        st.error(f"Preview Generation Failed: {e}") 

# Function to process the extracted text, show the preview and analyze the document based on the extracted text
def process_extracted_text(document_text, file_type,selected_model, temperature, tone, rewrite_target, target_length):
    if document_text.strip():
        st.success(
            f"{file_type} extracted successfully ✅"
        )
        show_preview(document_text)
        analyze_document(document_text, selected_model, temperature, tone, rewrite_target, target_length)
    else:
        st.error(
            f"No text extracted from {file_type}"
        )

def app():
    # side bar 
    with st.sidebar:
        st.title('🦙💬 Llama summarizer')
        
        models = ollama.list()
    # To choose all the models present in the ollama list and show them in the dropdown menu
        model_names = [
            model.model
            for model in models.models
    ]
        selected_model = st.selectbox("Choose Model",model_names)


        # To choose the temperature and show it in the slider
        temperature = st.slider("Temperature",0.0,1.0,0.7)


        # To choose the summary length and show it in the dropdown menu
        st.subheader("Summary Settings")

        summary_size = st.selectbox(
            "Summary Length",
            ["Short", "Medium", "Long"]
        )
        summary_config = {
        "Short": "100-200 words",
        "Medium": "300-500 words",
        "Long": "700-1000 words"
    }
        target_length = summary_config[summary_size] # select max lenght based on words count  to avoid the summary being cut off due to token limits.



        # To choose the rewrite tone and show it in the dropdown menu
        tone = st.selectbox(
            "Rewrite Tone",
            [
                "Professional",
                "Formal",
                "Casual",
                "Executive",
                "Technical",
                "Academic",
                "Marketing",
                "Simple English"
            ]
        )

        rewrite_target = st.selectbox(
        "Rewrite",
        [
            "Summary",
            "Full Document"
        ]
    )
        




    # Main app
    st.title("📄 AI Document Analyzer")

    tab1, tab2 = st.tabs(
        ["📝 Text Document", "📂 Upload Document"]
    )

    with tab1:
    #Text Box to paste the document to be analyzed
        document_text = st.text_area(
        "Paste your text here:",
        height=300,
        placeholder="Enter the text you want to summarize or rewrite."
    )


        # Button to trigger the analysis
        analyze_button = st.button("Analyze Document")
        #analyze_button Logic
        if analyze_button:
            if not document_text.strip():
                st.warning("Please enter a document.")
            
            else:
                analyze_document(document_text,selected_model,temperature,tone,rewrite_target,target_length)


    #File uploader to upload the document to be analyzed 
    with tab2:
        uploaded_file = st.file_uploader(
            "Upload a document",
            type=["pdf", "docx", "pptx", "txt"]
        )
        # Button to trigger the analysis of the uploaded document
        analyze_upload = st.button(
            "Analyze Uploaded Document"
        )
        #Analyze_Upload Button Logic
        if analyze_upload:

            if not uploaded_file:

                st.warning("Please upload a document.")

            else:

                try:
                    st.success("Document uploaded successfully ✅")
                    with st.expander("📄 Document Details"):
                        st.write(f"📄 File Name: {uploaded_file.name}")
                        st.write(f"📁 File Type: {uploaded_file.type}")
                        st.write(
                            f"📊 Size: {uploaded_file.size / 1024:.2f} KB"
                        )

                    file_name = uploaded_file.name.lower()
                #For each file type, extract the text and then show the preview and analyze the document based on the extracted text.
                    if file_name.endswith(".txt"):#for .txt
                        document_text = (
                            uploaded_file
                            .read()
                            .decode("utf-8")
                        )
                        process_extracted_text(document_text, "Text", selected_model,temperature,tone,rewrite_target,target_length)
    


                    elif file_name.endswith(".pdf"): #for .pdf
                        document_text = extract_pdf(uploaded_file)
                        process_extracted_text(document_text, "PDF", selected_model,temperature,tone,rewrite_target,target_length)
                                                                                    
                        


                    elif file_name.endswith(".docx"):#for .docx
                        document_text = extract_docx(uploaded_file)
                        process_extracted_text(document_text, "DOCX", selected_model,temperature,tone,rewrite_target,target_length)
                                                                                    


                    elif file_name.endswith(".pptx"):#for .pptx
                        document_text = extract_pptx(uploaded_file)
                        process_extracted_text(document_text, "PPTX", selected_model,temperature,tone,rewrite_target,target_length)
                                                                                    


                    else:

                        st.error(
                            "Unsupported file type."
                        )

                except Exception as e:

                    st.error(
                        f"Error reading file: {e}"
                    )



if __name__ == "__main__":
    app()